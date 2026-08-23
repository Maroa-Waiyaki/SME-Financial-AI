"""Document parser for multi-format SME business documents.

Supports:
- PDF (.pdf)
- Microsoft Word (.docx, .doc)
- Microsoft Excel (.xlsx, .xls)
- Microsoft PowerPoint (.pptx, .ppt)
- Markdown & Plain text (.md, .txt, .csv)
"""

from __future__ import annotations

import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """Extract clean text content from various file formats."""
    ext = Path(filename).suffix.lower()

    if ext in {".txt", ".md", ".csv", ".json"}:
        try:
            return file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            return file_bytes.decode("latin-1", errors="ignore")

    if ext == ".pdf":
        return _extract_pdf(file_bytes)

    if ext == ".docx":
        return _extract_docx(file_bytes)

    if ext in {".xlsx", ".xls"}:
        return _extract_excel(file_bytes)

    if ext in {".pptx", ".ppt"}:
        return _extract_pptx(file_bytes)

    # Fallback to string decode
    try:
        return file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        return ""


def _extract_pdf(file_bytes: bytes) -> str:
    try:
        import pypdf

        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        pages_text = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                pages_text.append(f"## Page {i + 1}\n{text.strip()}")
        return "\n\n".join(pages_text)
    except Exception as exc:
        logger.error("Failed to parse PDF: %s", exc)
        return ""


def _extract_docx(file_bytes: bytes) -> str:
    try:
        import docx

        doc = docx.Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    paragraphs.append(" | ".join(row_text))
        return "\n\n".join(paragraphs)
    except Exception as exc:
        logger.error("Failed to parse DOCX: %s", exc)
        return ""


def _extract_excel(file_bytes: bytes) -> str:
    try:
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = [f"## Sheet: {sheet_name}"]
            for row in ws.iter_rows(values_only=True):
                non_empty = [str(val).strip() for val in row if val is not None and str(val).strip()]
                if non_empty:
                    rows_text.append(" | ".join(non_empty))
            if len(rows_text) > 1:
                sheets_text.append("\n".join(rows_text))
        return "\n\n".join(sheets_text)
    except Exception as exc:
        logger.error("Failed to parse Excel: %s", exc)
        return ""


def _extract_pptx(file_bytes: bytes) -> str:
    try:
        import pptx

        prs = pptx.Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_content = [f"## Slide {i + 1}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_content.append(text)
            if len(slide_content) > 1:
                slides_text.append("\n".join(slide_content))
        return "\n\n".join(slides_text)
    except Exception as exc:
        logger.error("Failed to parse PPTX: %s", exc)
        return ""
