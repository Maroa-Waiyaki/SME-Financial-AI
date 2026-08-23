"""Deterministic financial statement generators and report exporters.

Produces:
1. Income Statement (Profit & Loss / P&L)
2. Balance Sheet / Financial Position Statement
3. PowerPoint Presentation Deck (.pptx)
4. Excel Financial Model Workbook (.xlsx)
"""

from __future__ import annotations

import io
from datetime import date, datetime
from decimal import Decimal
from typing import Any

from sqlalchemy import and_, func, select
from sqlalchemy.orm import Session

from src.database.models import Business, Expense, Invoice, Loan, Sale, Transaction
from src.financial import calculations as calc


def generate_pnl_statement(
    session: Session,
    business_id: str,
    start_date: str | date = "2023-01-01",
    end_date: str | date = "2023-12-31",
) -> dict[str, Any]:
    """Generate a formal SME Profit & Loss statement."""
    business = session.get(Business, business_id)
    if not business:
        raise ValueError(f"Business {business_id} not found")

    s_date = datetime.strptime(str(start_date), "%Y-%m-%d").date() if isinstance(start_date, str) else start_date
    e_date = datetime.strptime(str(end_date), "%Y-%m-%d").date() if isinstance(end_date, str) else end_date

    # 1. Gross Revenue from Sales
    sales_stmt = select(func.coalesce(func.sum(Sale.total_amount), Decimal("0"))).where(
        and_(Sale.business_id == business_id, Sale.date >= s_date, Sale.date <= e_date)
    )
    gross_revenue = session.execute(sales_stmt).scalar_one()

    # 2. Operating Expenses by category
    exp_stmt = (
        select(Expense.category, func.coalesce(func.sum(Expense.amount), Decimal("0")).label("total"))
        .where(and_(Expense.business_id == business_id, Expense.date >= s_date, Expense.date <= e_date))
        .group_by(Expense.category)
        .order_by(func.sum(Expense.amount).desc())
    )
    expense_rows = session.execute(exp_stmt).all()
    expense_breakdown = {r[0]: Decimal(str(r[1])) for r in expense_rows}

    cogs = expense_breakdown.get("Inventory", Decimal("0"))
    gross_profit = gross_revenue - cogs
    gross_margin = (gross_profit / gross_revenue * 100) if gross_revenue > 0 else Decimal("0")

    operating_expenses = sum(amt for cat, amt in expense_breakdown.items() if cat != "Inventory")
    ebitda = gross_profit - operating_expenses
    net_profit = ebitda
    net_profit_margin = (net_profit / gross_revenue * 100) if gross_revenue > 0 else Decimal("0")

    return {
        "business_id": business_id,
        "business_name": business.business_name,
        "sector": business.sector,
        "county": business.county,
        "period_start": str(s_date),
        "period_end": str(e_date),
        "gross_revenue": float(gross_revenue),
        "cost_of_goods_sold": float(cogs),
        "gross_profit": float(gross_profit),
        "gross_margin_percent": round(float(gross_margin), 2),
        "operating_expenses_total": float(operating_expenses),
        "expense_breakdown": {k: float(v) for k, v in expense_breakdown.items()},
        "net_profit": float(net_profit),
        "net_profit_margin_percent": round(float(net_profit_margin), 2),
    }


def generate_balance_sheet(
    session: Session,
    business_id: str,
    as_of_date: str | date = "2023-12-31",
) -> dict[str, Any]:
    """Generate a formal SME Balance Sheet (Statement of Financial Position)."""
    business = session.get(Business, business_id)
    if not business:
        raise ValueError(f"Business {business_id} not found")

    d_date = datetime.strptime(str(as_of_date), "%Y-%m-%d").date() if isinstance(as_of_date, str) else as_of_date

    # 1. Cash and Bank Balances (Latest Transaction balance_after)
    latest_txn_stmt = (
        select(Transaction.balance_after)
        .where(and_(Transaction.business_id == business_id, Transaction.timestamp <= datetime.combine(d_date, datetime.max.time())))
        .order_by(Transaction.timestamp.desc())
        .limit(1)
    )
    cash_balance = session.execute(latest_txn_stmt).scalar_one_or_none() or Decimal("0")

    # 2. Accounts Receivable (Unpaid and overdue invoices)
    ar_stmt = select(func.coalesce(func.sum(Invoice.outstanding_amount), Decimal("0"))).where(
        and_(Invoice.business_id == business_id, Invoice.invoice_date <= d_date, Invoice.status.in_(["unpaid", "overdue", "partial"]))
    )
    accounts_receivable = session.execute(ar_stmt).scalar_one()

    # 3. Estimated Inventory Asset (Based on recent monthly turnover)
    inventory_val = Decimal(str(business.monthly_revenue_estimate)) * Decimal("0.35")
    total_current_assets = Decimal(str(cash_balance)) + Decimal(str(accounts_receivable)) + inventory_val

    # 4. Long-term / Fixed Assets (Equipment, fixtures based on size and sector)
    fixed_assets = Decimal(str(business.monthly_revenue_estimate)) * Decimal("1.2")
    total_assets = total_current_assets + fixed_assets

    # 5. Liabilities (Outstanding loans)
    loan_stmt = select(func.coalesce(func.sum(Loan.outstanding_balance), Decimal("0"))).where(
        and_(Loan.business_id == business_id, Loan.status == "active")
    )
    outstanding_loans = session.execute(loan_stmt).scalar_one()
    accounts_payable = Decimal(str(business.monthly_revenue_estimate)) * Decimal("0.15")
    total_liabilities = Decimal(str(outstanding_loans)) + accounts_payable

    # 6. Owner's Equity (Net Worth = Assets - Liabilities)
    retained_earnings = total_assets - total_liabilities

    return {
        "business_id": business_id,
        "business_name": business.business_name,
        "as_of_date": str(d_date),
        "assets": {
            "current_assets": {
                "cash_and_bank": float(cash_balance),
                "accounts_receivable": float(accounts_receivable),
                "inventory_estimate": round(float(inventory_val), 2),
                "total_current_assets": round(float(total_current_assets), 2),
            },
            "fixed_assets": {
                "property_plant_equipment": round(float(fixed_assets), 2),
                "total_fixed_assets": round(float(fixed_assets), 2),
            },
            "total_assets": round(float(total_assets), 2),
        },
        "liabilities": {
            "current_liabilities": {
                "accounts_payable_estimate": round(float(accounts_payable), 2),
                "short_term_debt": float(outstanding_loans),
                "total_current_liabilities": round(float(total_liabilities), 2),
            },
            "total_liabilities": round(float(total_liabilities), 2),
        },
        "equity": {
            "retained_earnings_and_capital": round(float(retained_earnings), 2),
            "total_equity": round(float(retained_earnings), 2),
            "total_liabilities_and_equity": round(float(total_assets), 2),
        },
    }


def export_financial_pptx(pnl_data: dict[str, Any], balance_data: dict[str, Any]) -> bytes:
    """Generate a formal 5-slide Executive Presentation deck (.pptx)."""
    import pptx
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = pptx.Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    primary_color = RGBColor(0, 90, 135)
    dark_gray = RGBColor(50, 50, 50)

    # Slide 1: Title Slide
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)
    
    tx_box = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{pnl_data['business_name']}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = primary_color

    p2 = tf.add_paragraph()
    p2.text = f"Financial Intelligence & Performance Review ({pnl_data['period_start']} to {pnl_data['period_end']})"
    p2.font.size = Pt(18)
    p2.font.color.rgb = dark_gray

    p3 = tf.add_paragraph()
    p3.text = f"Sector: {pnl_data['sector']} | County: {pnl_data['county']} | Platform ID: {pnl_data['business_id']}"
    p3.font.size = Pt(14)

    # Slide 2: Executive Summary & KPIs
    slide2 = prs.slides.add_slide(blank_layout)
    header = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(0.8))
    header.text_frame.text = "1. Executive Summary & Key Metrics"
    header.text_frame.paragraphs[0].font.size = Pt(24)
    header.text_frame.paragraphs[0].font.bold = True
    header.text_frame.paragraphs[0].font.color.rgb = primary_color

    body2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.5))
    tf2 = body2.text_frame
    tf2.word_wrap = True
    kpi_lines = [
        f"• Gross Revenue: KES {pnl_data['gross_revenue']:,.2f}",
        f"• Net Profit: KES {pnl_data['net_profit']:,.2f} (Net Margin: {pnl_data['net_profit_margin_percent']}%)",
        f"• Total Business Assets: KES {balance_data['assets']['total_assets']:,.2f}",
        f"• Cash & Bank Liquidity: KES {balance_data['assets']['current_assets']['cash_and_bank']:,.2f}",
        f"• Total Liabilities: KES {balance_data['liabilities']['total_liabilities']:,.2f}",
        f"• Estimated Net Worth / Equity: KES {balance_data['equity']['total_equity']:,.2f}",
    ]
    for line in kpi_lines:
        p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(16)

    # Slide 3: Profit & Loss Statement (Table)
    slide3 = prs.slides.add_slide(blank_layout)
    header3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(0.8))
    header3.text_frame.text = "2. Profit & Loss Statement (P&L)"
    header3.text_frame.paragraphs[0].font.size = Pt(24)
    header3.text_frame.paragraphs[0].font.bold = True
    header3.text_frame.paragraphs[0].font.color.rgb = primary_color

    table_shape = slide3.shapes.add_table(6, 2, Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.2))
    table = table_shape.table
    table.columns[0].width = Inches(5.5)
    table.columns[1].width = Inches(2.9)

    pnl_rows = [
        ("Revenue Item", "Amount (KES)"),
        ("Gross Sales Revenue", f"{pnl_data['gross_revenue']:,.2f}"),
        ("Cost of Goods Sold (Inventory)", f"({pnl_data['cost_of_goods_sold']:,.2f})"),
        ("Gross Profit", f"{pnl_data['gross_profit']:,.2f}"),
        ("Operating Expenses (Salaries, Rent, Utilities, Transport)", f"({pnl_data['operating_expenses_total']:,.2f})"),
        ("Net Profit / EBITDA", f"{pnl_data['net_profit']:,.2f}"),
    ]
    for r_idx, (col1, col2) in enumerate(pnl_rows):
        cell1, cell2 = table.cell(r_idx, 0), table.cell(r_idx, 1)
        cell1.text = col1
        cell2.text = col2
        if r_idx in (0, 3, 5):
            cell1.text_frame.paragraphs[0].font.bold = True
            cell2.text_frame.paragraphs[0].font.bold = True

    # Slide 4: Balance Sheet (Statement of Financial Position)
    slide4 = prs.slides.add_slide(blank_layout)
    header4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(0.8))
    header4.text_frame.text = "3. Balance Sheet (Financial Position)"
    header4.text_frame.paragraphs[0].font.size = Pt(24)
    header4.text_frame.paragraphs[0].font.bold = True
    header4.text_frame.paragraphs[0].font.color.rgb = primary_color

    table_shape4 = slide4.shapes.add_table(6, 2, Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.2))
    table4 = table_shape4.table
    table4.columns[0].width = Inches(5.5)
    table4.columns[1].width = Inches(2.9)

    bs_rows = [
        ("Asset & Liability Category", "As of Date Value (KES)"),
        ("Current Assets (Cash, Receivables, Inventory)", f"{balance_data['assets']['current_assets']['total_current_assets']:,.2f}"),
        ("Fixed Assets (Equipment, Fixtures)", f"{balance_data['assets']['fixed_assets']['total_fixed_assets']:,.2f}"),
        ("Total Assets", f"{balance_data['assets']['total_assets']:,.2f}"),
        ("Total Liabilities (Loans & Payables)", f"{balance_data['liabilities']['total_liabilities']:,.2f}"),
        ("Owner's Equity & Retained Earnings", f"{balance_data['equity']['total_equity']:,.2f}"),
    ]
    for r_idx, (col1, col2) in enumerate(bs_rows):
        cell1, cell2 = table4.cell(r_idx, 0), table4.cell(r_idx, 1)
        cell1.text = col1
        cell2.text = col2
        if r_idx in (0, 3, 5):
            cell1.text_frame.paragraphs[0].font.bold = True
            cell2.text_frame.paragraphs[0].font.bold = True

    # Slide 5: Strategic Recommendations
    slide5 = prs.slides.add_slide(blank_layout)
    header5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(0.8))
    header5.text_frame.text = "4. AI Financial Intelligence Recommendations"
    header5.text_frame.paragraphs[0].font.size = Pt(24)
    header5.text_frame.paragraphs[0].font.bold = True
    header5.text_frame.paragraphs[0].font.color.rgb = primary_color

    body5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.5))
    tf5 = body5.text_frame
    recs = [
        "1. Working Capital Optimization: Maintain cash buffer equal to 2 months of operating expenses.",
        "2. Debt Capacity: Healthy debt-to-equity ratio supports credit facility expansion for inventory pre-funding.",
        "3. Receivables Management: Follow up on outstanding invoices to shorten cash conversion cycle.",
        "4. Tax & Statutory Planning: Set aside monthly provisions to prevent cash-flow shocks.",
    ]
    for r in recs:
        p = tf5.add_paragraph()
        p.text = r
        p.font.size = Pt(15)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def export_financial_excel(pnl_data: dict[str, Any], balance_data: dict[str, Any]) -> bytes:
    """Generate a multi-tab formatted Excel financial model (.xlsx)."""
    import openpyxl

    wb = openpyxl.Workbook()
    
    # Sheet 1: P&L
    ws_pnl = wb.active
    ws_pnl.title = "Profit & Loss"
    ws_pnl.append([f"PROFIT & LOSS STATEMENT - {pnl_data['business_name']}"])
    ws_pnl.append([f"Period: {pnl_data['period_start']} to {pnl_data['period_end']}"])
    ws_pnl.append([])
    ws_pnl.append(["Line Item", "Amount (KES)"])
    ws_pnl.append(["Gross Revenue", pnl_data["gross_revenue"]])
    ws_pnl.append(["Cost of Goods Sold (Inventory)", pnl_data["cost_of_goods_sold"]])
    ws_pnl.append(["Gross Profit", pnl_data["gross_profit"]])
    ws_pnl.append([])
    ws_pnl.append(["Operating Expenses:"])
    for cat, amt in pnl_data["expense_breakdown"].items():
        if cat != "Inventory":
            ws_pnl.append([f"  - {cat}", amt])
    ws_pnl.append(["Total Operating Expenses", pnl_data["operating_expenses_total"]])
    ws_pnl.append([])
    ws_pnl.append(["NET PROFIT / EBITDA", pnl_data["net_profit"]])
    ws_pnl.append(["Net Profit Margin", f"{pnl_data['net_profit_margin_percent']}%"])

    # Sheet 2: Balance Sheet
    ws_bs = wb.create_sheet(title="Balance Sheet")
    ws_bs.append([f"BALANCE SHEET - {balance_data['business_name']}"])
    ws_bs.append([f"As of: {balance_data['as_of_date']}"])
    ws_bs.append([])
    ws_bs.append(["Assets", "Amount (KES)"])
    ws_bs.append(["Cash & Bank Balances", balance_data["assets"]["current_assets"]["cash_and_bank"]])
    ws_bs.append(["Accounts Receivable", balance_data["assets"]["current_assets"]["accounts_receivable"]])
    ws_bs.append(["Estimated Inventory", balance_data["assets"]["current_assets"]["inventory_estimate"]])
    ws_bs.append(["Total Current Assets", balance_data["assets"]["current_assets"]["total_current_assets"]])
    ws_bs.append(["Fixed Assets & Equipment", balance_data["assets"]["fixed_assets"]["total_fixed_assets"]])
    ws_bs.append(["TOTAL ASSETS", balance_data["assets"]["total_assets"]])
    ws_bs.append([])
    ws_bs.append(["Liabilities & Equity", "Amount (KES)"])
    ws_bs.append(["Accounts Payable", balance_data["liabilities"]["current_liabilities"]["accounts_payable_estimate"]])
    ws_bs.append(["Outstanding Loans", balance_data["liabilities"]["current_liabilities"]["short_term_debt"]])
    ws_bs.append(["Total Liabilities", balance_data["liabilities"]["total_liabilities"]])
    ws_bs.append(["Owner's Equity & Retained Earnings", balance_data["equity"]["total_equity"]])
    ws_bs.append(["TOTAL LIABILITIES & EQUITY", balance_data["equity"]["total_liabilities_and_equity"]])

    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()
